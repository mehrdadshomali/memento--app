import pptx

prs = pptx.Presentation("/Users/mehrdadshomali/Desktop/NKU_BMU_Poster_Sablonu_A3_Dikey (2).pptx")

def replace_text(shape, new_text):
    if hasattr(shape, "text_frame"):
        # preserve paragraph styling by only replacing text of the first run
        # but since we might have multi-line text, we can just assign text. 
        # assigning text clears formatting but keeps the default shape format.
        shape.text = new_text

slide = prs.slides[0]

# Shape 7: Poster Title
slide.shapes[7].text = "MEMENTO: Alzheimer Hastaları İçin Terapötik Hafıza Egzersiz Mobil Uygulaması"

# Shape 8: Authors
slide.shapes[8].text = "Mehrdad Shomali (2210656637)  ·  Kerim Can Güzel (2220656047)  ·  Efekan Çetin (222656034)"

# Shape 9: Advisor
slide.shapes[9].text = "¹Tekirdağ Namık Kemal Üniversitesi\nDanışman Öğretim Üyesi: Erkan Özhan"

# Shape 13: ÖZET content
slide.shapes[13].text = "Bu poster, Alzheimer ve demans hastalarının bilişsel sağlığını desteklemek amacıyla geliştirilmiş 'Memento' mobil uygulamasının detaylarını ve çalışma prensiplerini sunmaktadır. React Native ve Expo teknolojileri kullanılarak iOS ve Android platformları için geliştirilen bu uygulama; hafıza oyunları, aile albümü, günlük rutin takibi, sesli mesajlar ve GPS tabanlı güvenlik modülü gibi kapsamlı özellikler içermektedir. Önerilen sistem, geleneksel yöntemlere kıyasla dijital bir terapi sunarak hastaların günlük aktivitelerini güvenle sürdürmelerini sağlamakta ve bakıcılara süreç takibi yapabilme imkanı vermektedir."

# Shape 17: GİRİŞ VE MOTİVASYON content
slide.shapes[17].text = "• Dünya genelinde 55 milyon Alzheimer hastası, günlük yaşam aktivitelerinde, yüz/ses hatırlamada ve mekân algısında ciddi zorluklar çekmektedir.\n• Mevcut piyasadaki sağlık uygulamalarının karmaşık arayüzleri (UI/UX), bilişsel bozukluğu olan hastalar için ciddi erişilebilirlik sorunları yaratmaktadır.\n• Amacımız: Hastaların sevdiklerini hatırlamasına destek olan, günlük ilaç/yemek rutinlerini takip eden ve evden uzaklaşıldığında kaybolma riskine karşı konum tabanlı güvenlik sağlayan hepsi bir arada bir mHealth (mobil sağlık) çözümü geliştirmek."

# Shape 20 & 21: İLGİLİ ÇALIŞMALAR -> YAZILIM MİMARİSİ
slide.shapes[20].text = "⚙️  YAZILIM MİMARİSİ VE TEKNOLOJİLER"
slide.shapes[21].text = "- Kullanıcı Arayüzü (UI): React Native (v0.81), Expo SDK (v54), TypeScript. WCAG AA standartlarına uygun, büyük dokunma hedefleri.\n- Durum Yönetimi (State): Context API ve Hooks ile modüler Provider mimarisi.\n- Konum ve Güvenlik: Expo Location ile arka planda GPS takibi. Mesafe hesaplamasında Haversine formülü kullanılmıştır.\n- Veritabanı: Çevrimdışı önbellekleme için AsyncStorage, bulut senkronizasyonu için Supabase."

# Shape 24 & 25: VERİ KÜMESİ -> MODÜLLER 1
slide.shapes[24].text = "📱  UYGULAMA MODÜLLERİ (1)"
slide.shapes[25].text = "1. Hafıza Oyunları Motoru: Fisher-Yates karıştırma algoritması kullanılarak dinamik olarak üretilen görsel (yüz) ve işitsel (ses) tanıma egzersizleri.\n2. Güvenlik ve Konum Takibi: Hasta belirlenen ev konumundan 100 metreden fazla uzaklaştığında, arka planda çalışan servis devreye girer ve hastaya eve dönmesi için otomatik bildirim gönderir."

# Shape 28 & 30: Görsel placeholder 1
slide.shapes[28].text = "📊  UYGULAMA ARAYÜZÜ - 1"
slide.shapes[30].text = "[ Buraya Uygulamanın Ana Ekranı ve Hafıza Oyunu Ekran Görüntülerini Ekleyebilirsiniz ]"

# Shape 33 & 34: DENEYSEL KURULUM -> MODÜLLER 2
slide.shapes[33].text = "📱  UYGULAMA MODÜLLERİ (2)"
slide.shapes[34].text = "3. Günlük Rutin Yönetimi: İlaç, yemek ve randevu hatırlatıcıları zamanlanarak hastaya bildirilir.\n4. Aile Albümü ve Sesli Mesajlar: Hastanın ailesiyle duygusal bağını korumak için tasarlanmış, .m4a ses kaydı destekli medya modülü.\n5. Bakıcı Modu ve Analiz: Hastanın rutin tamamlama istatistiklerini ve oyun skorlarını gösteren yetkilendirilmiş ekran."

# Shape 37 & 38: ÖNERİLEN YÖNTEM -> AKIŞ DİYAGRAMI
slide.shapes[37].text = "⚙️  SİSTEM AKIŞ DİYAGRAMI"
slide.shapes[38].text = "[ Buraya Tez Raporunuzdaki Uygulama Akış Şeması veya Mimari Çizimini Ekleyebilirsiniz ]"

# Shape 41 & 47: SONUÇ GRAFİKLERİ -> Görsel placeholder 2
slide.shapes[41].text = "📈  UYGULAMA ARAYÜZÜ - 2"
slide.shapes[47].text = "[ Buraya Aile Albümü, Bakıcı Ekranı ve Konum Haritası Ekran Görüntülerini Ekleyebilirsiniz ]"

# Shape 50 & 51: TARTIŞMA
slide.shapes[50].text = "💬  TARTIŞMA"
slide.shapes[51].text = "• Memento, basit ve yönlendirici kullanıcı arayüzü sayesinde yaşlı ve bilişsel yetisi kısıtlı bireyler tarafından stres olmadan kolayca kullanılabilir.\n• Haversine formülü tabanlı konum hesaplaması ve arka plan konum servisi (Background Task), yüksek GPS doğruluğu sağlarken cihazın pil tüketimini optimize edecek şekilde tasarlanmıştır."

# Shape 54 & 55: SONUÇ VE GELECEK ÇALIŞMA
slide.shapes[54].text = "✅  SONUÇ VE GELECEK ÇALIŞMA"
slide.shapes[55].text = "Sonuç: Memento uygulaması, Alzheimer hastalarının bağımsızlıklarını artırmayı ve bakıcıların üzerindeki yükü hafifletmeyi başaran fonksiyonel bir mobil sağlık sistemidir.\n\nGelecek Çalışmalar:\n• Akıllı saatler ile kalp atışı ve düşme sensörü entegrasyonu.\n• Makine öğrenmesi algoritmaları ile hastanın günlük duygu durumu analizi."

# Shape 58 & 59: İletişim
slide.shapes[58].text = "📧 mehrdad.shomali@nku.edu.tr  |  📍 Çorlu Mühendislik Fakültesi, Tekirdağ"
slide.shapes[59].text = "Bu proje çalışması Tekirdağ Namık Kemal Üniversitesi Bilgisayar Mühendisliği bölümünde gerçekleştirilmiştir."

prs.save("/Users/mehrdadshomali/Desktop/Memento_Poster.pptx")
print("Poster başarıyla oluşturuldu!")
